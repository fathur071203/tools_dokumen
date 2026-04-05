from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import fitz
from PIL import Image
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Emu
from pdf2docx import Converter

from src.services.output_naming_service import OutputNamingService


class ConvertService:
    PDF_TEXT_MIN_FOR_STRUCTURED_DOCX = 50
    PDF_TEXT_MIN_FOR_IMAGE_HEAVY_OVERRIDE = 500
    PDF_IMAGE_HEAVY_TOTAL_IMAGES = 10
    PDF_PAGE_TEXT_MIN = 80
    PDF_PAGE_IMAGE_HEAVY_MIN_IMAGES = 5
    PDF_PAGE_IMAGE_HEAVY_MAX_TEXT = 300

    SUPPORTED_OPTIONS: dict[str, list[str]] = {
        ".pdf": ["docx", "pptx", "xlsx"],
        ".doc": ["pdf"],
        ".docx": ["pdf"],
        ".xls": ["pdf"],
        ".xlsx": ["pdf"],
        ".ppt": ["pdf"],
        ".pptx": ["pdf"],
        ".png": ["pdf"],
        ".jpg": ["pdf"],
        ".jpeg": ["pdf"],
        ".bmp": ["pdf"],
        ".webp": ["pdf"],
    }

    MIME_BY_EXT = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }

    @staticmethod
    def supports_office_to_pdf() -> bool:
        return sys.platform.startswith("win")

    @classmethod
    def get_available_targets(cls, filename: str) -> list[str]:
        source_ext = Path(filename).suffix.lower()
        targets = list(cls.SUPPORTED_OPTIONS.get(source_ext, []))
        if source_ext in {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"} and not cls.supports_office_to_pdf():
            return []
        return targets

    @classmethod
    def get_available_targets_for_uploads(cls, uploads: list[Any]) -> list[str]:
        if not uploads:
            return []

        if len(uploads) == 1:
            return cls.get_available_targets(uploads[0].name)

        if any(Path(getattr(upload, "name", "")).suffix.lower() == ".pdf" for upload in uploads):
            return []

        common_targets = set(cls.get_available_targets(uploads[0].name))
        for upload in uploads[1:]:
            common_targets &= set(cls.get_available_targets(upload.name))

        if "pdf" in common_targets:
            return ["pdf"]
        return []

    @classmethod
    def convert_files(
        cls,
        uploads: list[Any],
        target_format: str,
        pdf_output_mode: str = "merge",
    ) -> tuple[io.BytesIO, str, str, str]:
        if not uploads:
            raise ValueError("Tidak ada file untuk dikonversi")

        if len(uploads) == 1:
            return cls.convert_file(uploads[0], target_format)

        target_format = target_format.lower()
        if target_format != "pdf":
            raise ValueError("Konversi multi-file saat ini hanya mendukung output PDF")
        if pdf_output_mode not in {"merge", "separate"}:
            raise ValueError("Mode output PDF tidak valid")

        outputs: list[tuple[str, bytes]] = []
        for upload in uploads:
            output, filename, _, _ = cls.convert_file(upload, target_format)
            outputs.append((filename, output.getvalue()))

        if pdf_output_mode == "merge":
            return cls._merge_pdf_outputs(outputs)
        return cls._package_pdf_outputs(outputs)

    @classmethod
    def convert_file(cls, upload: Any, target_format: str) -> tuple[io.BytesIO, str, str, str]:
        source_name = getattr(upload, "name", "document")
        source_ext = Path(source_name).suffix.lower()
        target_format = target_format.lower()

        if target_format not in cls.get_available_targets(source_name):
            raise ValueError(f"Konversi dari {source_ext or 'file ini'} ke {target_format.upper()} belum didukung")

        upload.seek(0)
        file_bytes = upload.read()

        if source_ext == ".pdf" and target_format == "docx":
            return cls._pdf_to_docx(file_bytes, source_name)
        if source_ext == ".pdf" and target_format == "pptx":
            return cls._pdf_to_pptx(file_bytes, source_name)
        if source_ext == ".pdf" and target_format == "xlsx":
            return cls._pdf_to_excel(file_bytes)
        if target_format == "pdf":
            return cls._convert_to_pdf(file_bytes, source_name)

        raise ValueError("Kombinasi konversi belum tersedia")

    @classmethod
    def _convert_to_pdf(cls, file_bytes: bytes, source_name: str) -> tuple[io.BytesIO, str, str, str]:
        source_ext = Path(source_name).suffix.lower()
        if source_ext in {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"}:
            return cls._office_to_pdf(file_bytes, source_name)
        if source_ext in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
            return cls._image_to_pdf(file_bytes, source_name)
        raise ValueError(f"Konversi dari {source_ext or 'file ini'} ke PDF belum didukung")

    @classmethod
    def _pdf_to_docx(cls, file_bytes: bytes, source_name: str) -> tuple[io.BytesIO, str, str, str]:
        processed_pdf = cls._preprocess_pdf(file_bytes)
        page_modes = cls._detect_pdf_page_modes(processed_pdf)
        total_pages = len(page_modes)
        image_pages = sum(1 for mode in page_modes if mode == "image")

        if total_pages and image_pages == total_pages:
            output = cls._pdf_to_docx_image_fallback(processed_pdf)
            return (
                output,
                OutputNamingService.build_filename("converted_document", ".docx"),
                cls.MIME_BY_EXT["docx"],
                "PDF terdeteksi dominan visual (design/scanned). Dikonversi ke Word mode visual (image-based) agar tata letak tetap mirip aslinya.",
            )

        if total_pages and 0 < image_pages < total_pages:
            output = cls._pdf_to_docx_page_hybrid(processed_pdf, page_modes)
            return (
                output,
                OutputNamingService.build_filename("converted_document", ".docx"),
                cls.MIME_BY_EXT["docx"],
                f"PDF dikonversi dengan Smart Hybrid per halaman: {total_pages - image_pages} halaman teks + {image_pages} halaman visual.",
            )

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            pdf_path = temp_path / source_name
            docx_path = temp_path / f"{Path(source_name).stem}.docx"
            pdf_path.write_bytes(processed_pdf)

            converter = Converter(str(pdf_path))
            try:
                try:
                    converter.convert(
                        str(docx_path),
                        start=0,
                        end=None,
                        layout=True,
                        multi_processing=True,
                    )
                except TypeError:
                    # Kompatibilitas jika versi pdf2docx tidak mendukung parameter tambahan.
                    converter.convert(str(docx_path))
            finally:
                converter.close()

            output = io.BytesIO(docx_path.read_bytes())
            output.seek(0)
            return (
                output,
                OutputNamingService.build_filename("converted_document", ".docx"),
                cls.MIME_BY_EXT["docx"],
                "PDF berhasil dikonversi ke Word (DOCX) dengan mode struktur dokumen.",
            )

    @classmethod
    def _preprocess_pdf(cls, file_bytes: bytes) -> bytes:
        """Bersihkan konten PDF untuk mengurangi noise layout sebelum konversi."""
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            for page in pdf:
                try:
                    page.clean_contents()
                except Exception:
                    continue
            return pdf.tobytes(garbage=3, deflate=True, clean=True)
        finally:
            pdf.close()

    @classmethod
    def _is_complex_layout(cls, file_bytes: bytes) -> bool:
        """Kompatibilitas API lama: dianggap kompleks jika tipe PDF bukan text-dominant."""
        return cls._detect_pdf_type(file_bytes) != "text"

    @classmethod
    def _classify_page_mode(cls, page: fitz.Page) -> str:
        """Klasifikasi mode per halaman: text vs image (design-heavy/scanned)."""
        text = (page.get_text("text") or "").strip()
        text_len = len(text)
        images = len(page.get_images(full=True))

        if text_len < cls.PDF_PAGE_TEXT_MIN:
            return "image"

        if images >= cls.PDF_PAGE_IMAGE_HEAVY_MIN_IMAGES and text_len < cls.PDF_PAGE_IMAGE_HEAVY_MAX_TEXT:
            return "image"

        return "text"

    @classmethod
    def _extract_clean_text_from_blocks(cls, page: fitz.Page) -> str:
        """Ekstrak teks berbasis blok agar urutan baca lebih natural dibanding mode text flatten."""
        blocks = page.get_text("blocks") or []
        blocks_sorted = sorted(blocks, key=lambda block: (block[1], block[0]))

        pieces: list[str] = []
        for block in blocks_sorted:
            text = str(block[4] or "").strip()
            if text:
                pieces.append(text)

        return "\n\n".join(pieces)

    @classmethod
    def _smart_paragraphs(cls, text: str) -> list[str]:
        """Bangun paragraf lebih rapi dari hasil ekstraksi blok."""
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not lines:
            return []

        paragraphs: list[str] = []
        current_parts: list[str] = []

        for line in lines:
            bullet_like = line.startswith(("-", "*", "•"))
            short_title_like = len(line) < 40

            if bullet_like or short_title_like:
                if current_parts:
                    paragraphs.append(" ".join(current_parts).strip())
                    current_parts = []
                paragraphs.append(line)
                continue

            current_parts.append(line)

        if current_parts:
            paragraphs.append(" ".join(current_parts).strip())

        return [paragraph for paragraph in paragraphs if paragraph]

    @classmethod
    def _apply_word_style(cls, paragraph, text: str) -> None:
        """Pilih style sederhana agar hasil Word tidak flat."""
        clean = (text or "").strip()
        if not clean:
            return

        if clean.isupper() and len(clean) <= 80:
            paragraph.style = "Heading 1"
            return

        if len(clean) <= 50:
            paragraph.style = "Heading 2"

    @classmethod
    def _detect_pdf_page_modes(cls, file_bytes: bytes) -> list[str]:
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            return [cls._classify_page_mode(page) for page in pdf]
        finally:
            pdf.close()

    @classmethod
    def _detect_pdf_type(cls, file_bytes: bytes) -> str:
        """
        Klasifikasi sederhana namun stabil:
        - text: dominan teks, cocok untuk pdf2docx
        - image-heavy: gambar banyak dan teks relatif sedikit
        - scanned: hampir tidak ada teks
        """
        page_modes = cls._detect_pdf_page_modes(file_bytes)
        total_pages = len(page_modes)
        image_pages = sum(1 for mode in page_modes if mode == "image")

        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        total_text_length = 0
        total_images = 0
        try:
            for page in pdf:
                text = page.get_text("text") or ""
                total_text_length += len(text.strip())
                total_images += len(page.get_images(full=True))
        finally:
            pdf.close()

        if total_text_length < cls.PDF_TEXT_MIN_FOR_STRUCTURED_DOCX:
            return "scanned"

        if total_pages and image_pages == total_pages:
            return "image-heavy"

        if (
            total_images > cls.PDF_IMAGE_HEAVY_TOTAL_IMAGES
            and total_text_length < cls.PDF_TEXT_MIN_FOR_IMAGE_HEAVY_OVERRIDE
        ):
            return "image-heavy"

        return "text"

    @classmethod
    def _pdf_to_docx_page_hybrid(cls, file_bytes: bytes, page_modes: list[str] | None = None) -> io.BytesIO:
        """Hybrid per halaman: halaman text jadi paragraf, halaman design/scanned jadi gambar."""
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        doc = Document()
        page_width_inches = 6.5

        if page_modes is None or len(page_modes) != len(pdf):
            page_modes = [cls._classify_page_mode(page) for page in pdf]

        try:
            for page_idx, page in enumerate(pdf):
                mode = page_modes[page_idx]

                if mode == "image":
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    image_stream = io.BytesIO(pix.tobytes("png"))
                    image_stream.seek(0)
                    doc.add_picture(image_stream, width=Inches(page_width_inches))
                else:
                    text = cls._extract_clean_text_from_blocks(page)
                    paragraphs = cls._smart_paragraphs(text)

                    if not paragraphs:
                        paragraphs = [text.strip()] if text.strip() else [""]

                    for paragraph_text in paragraphs:
                        paragraph = doc.add_paragraph(paragraph_text)
                        cls._apply_word_style(paragraph, paragraph_text)

                if page_idx < len(pdf) - 1:
                    doc.add_page_break()

            output = io.BytesIO()
            doc.save(output)
            output.seek(0)
            return output
        finally:
            pdf.close()

    @classmethod
    def _pdf_to_docx_image_fallback(cls, file_bytes: bytes) -> io.BytesIO:
        """Fallback visual-preserving: tiap halaman PDF dijadikan gambar di Word."""
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        doc = Document()
        page_width_inches = 6.5

        try:
            for page_idx, page in enumerate(pdf):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image_stream = io.BytesIO(pix.tobytes("png"))
                image_stream.seek(0)

                doc.add_picture(image_stream, width=Inches(page_width_inches))
                if page_idx < len(pdf) - 1:
                    doc.add_page_break()

            output = io.BytesIO()
            doc.save(output)
            output.seek(0)
            return output
        finally:
            pdf.close()

    @classmethod
    def _pdf_to_pptx(cls, file_bytes: bytes, source_name: str) -> tuple[io.BytesIO, str, str, str]:
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        presentation = Presentation()
        presentation.slide_width = Emu(9144000)
        presentation.slide_height = Emu(5143500)

        blank_layout = presentation.slide_layouts[6]

        for page in pdf:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_stream = io.BytesIO(pix.tobytes("png"))
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image_stream,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        if len(presentation.slides) > len(pdf):
            first_slide_id = presentation.slides._sldIdLst[0]
            presentation.slides._sldIdLst.remove(first_slide_id)

        output = io.BytesIO()
        presentation.save(output)
        output.seek(0)
        pdf.close()
        filename = OutputNamingService.build_filename("converted_document", ".pptx")
        return output, filename, cls.MIME_BY_EXT["pptx"], "PDF berhasil dikonversi ke PowerPoint (setiap halaman menjadi slide)."

    @classmethod
    def _image_to_pdf(cls, file_bytes: bytes, source_name: str) -> tuple[io.BytesIO, str, str, str]:
        with Image.open(io.BytesIO(file_bytes)) as img:
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            output = io.BytesIO()
            img.save(output, format="PDF", resolution=150.0)
            output.seek(0)
            filename = OutputNamingService.build_filename("converted_document", ".pdf")
            return output, filename, cls.MIME_BY_EXT["pdf"], "Gambar berhasil dikonversi ke PDF."

    @classmethod
    def _pdf_to_excel(cls, file_bytes: bytes) -> tuple[io.BytesIO, str, str, str]:
        import pandas as pd

        output = io.BytesIO()
        pdf = fitz.open(stream=file_bytes, filetype="pdf")

        try:
            with pd.ExcelWriter(output, engine="openpyxl") as writer:
                sheet_index = 1
                has_written_content = False

                for page_number, page in enumerate(pdf, start=1):
                    try:
                        table_finder = page.find_tables()
                        tables = table_finder.tables if table_finder else []
                    except Exception:
                        tables = []

                    for table in tables:
                        table_rows = table.extract() or []
                        cleaned_rows = [row for row in table_rows if any((cell or "").strip() for cell in row)]
                        if not cleaned_rows:
                            continue

                        header = [str(col).strip() if col is not None else "" for col in cleaned_rows[0]]
                        data_rows = cleaned_rows[1:] if len(cleaned_rows) > 1 else []
                        dataframe = pd.DataFrame(data_rows, columns=header)

                        sheet_name = f"Halaman_{page_number}_{sheet_index}"
                        dataframe.to_excel(writer, sheet_name=sheet_name[:31], index=False)
                        sheet_index += 1
                        has_written_content = True

                if not has_written_content:
                    fallback_rows: list[dict[str, str]] = []
                    for page_number, page in enumerate(pdf, start=1):
                        text = page.get_text("text") or ""
                        lines = [line.strip() for line in text.splitlines() if line.strip()]
                        for line_number, line in enumerate(lines, start=1):
                            fallback_rows.append(
                                {
                                    "halaman": str(page_number),
                                    "baris": str(line_number),
                                    "teks": line,
                                }
                            )

                    fallback_df = pd.DataFrame(fallback_rows or [{"halaman": "", "baris": "", "teks": ""}])
                    fallback_df.to_excel(writer, sheet_name="PDF_Text", index=False)
        finally:
            pdf.close()

        output.seek(0)
        filename = OutputNamingService.build_filename("converted_document", ".xlsx")
        return (
            output,
            filename,
            cls.MIME_BY_EXT["xlsx"],
            "PDF berhasil dikonversi ke Excel (tabel terdeteksi otomatis; jika tidak ada tabel, teks diekspor).",
        )

    @classmethod
    def _office_to_pdf(cls, file_bytes: bytes, source_name: str) -> tuple[io.BytesIO, str, str, str]:
        if not cls.supports_office_to_pdf():
            raise RuntimeError("Konversi Office ke PDF hanya tersedia di Windows yang memiliki Microsoft Office desktop.")

        try:
            import pythoncom
            import win32com.client
        except ImportError as exc:
            raise ImportError("Dependensi Windows untuk konversi Office belum tersedia.") from exc

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            source_path = temp_path / source_name
            output_path = temp_path / f"{Path(source_name).stem}.pdf"
            source_path.write_bytes(file_bytes)
            source_ext = source_path.suffix.lower()

            pythoncom.CoInitialize()
            try:
                if source_ext in {".doc", ".docx"}:
                    cls._word_to_pdf(source_path, output_path, win32com.client)
                elif source_ext in {".xls", ".xlsx"}:
                    cls._excel_to_pdf(source_path, output_path, win32com.client)
                elif source_ext in {".ppt", ".pptx"}:
                    cls._powerpoint_to_pdf(source_path, output_path, win32com.client)
                else:
                    raise ValueError("Format Office belum didukung")
            except Exception as exc:
                raise RuntimeError(
                    "Konversi Office ke PDF membutuhkan Microsoft Office desktop terpasang di Windows."
                ) from exc
            finally:
                pythoncom.CoUninitialize()

            output = io.BytesIO(output_path.read_bytes())
            output.seek(0)
            return output, OutputNamingService.build_filename("converted_document", ".pdf"), cls.MIME_BY_EXT["pdf"], "File Office berhasil dikonversi ke PDF."

    @classmethod
    def _merge_pdf_outputs(cls, outputs: list[tuple[str, bytes]]) -> tuple[io.BytesIO, str, str, str]:
        merged_pdf = fitz.open()
        try:
            for _, pdf_bytes in outputs:
                source_pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
                try:
                    merged_pdf.insert_pdf(source_pdf)
                finally:
                    source_pdf.close()

            output = io.BytesIO(merged_pdf.tobytes())
            output.seek(0)
            return (
                output,
                OutputNamingService.build_filename("converted_document", ".pdf"),
                cls.MIME_BY_EXT["pdf"],
                f"{len(outputs)} file berhasil dikonversi dan digabung menjadi 1 PDF.",
            )
        finally:
            merged_pdf.close()

    @classmethod
    def _package_pdf_outputs(cls, outputs: list[tuple[str, bytes]]) -> tuple[io.BytesIO, str, str, str]:
        outputs = OutputNamingService.anonymize_named_payloads(outputs, "converted_document", ".pdf")
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as archive:
            for filename, file_bytes in outputs:
                archive.writestr(filename, file_bytes)
        zip_buffer.seek(0)
        return (
            zip_buffer,
            OutputNamingService.build_filename("converted_documents", ".zip"),
            "application/zip",
            f"{len(outputs)} file berhasil dikonversi menjadi PDF terpisah.",
        )

    @staticmethod
    def _word_to_pdf(source_path: Path, output_path: Path, client_module: Any) -> None:
        word = client_module.DispatchEx("Word.Application")
        word.Visible = False
        document = None
        try:
            document = word.Documents.Open(str(source_path))
            document.SaveAs(str(output_path), FileFormat=17)
        finally:
            if document is not None:
                document.Close(False)
            word.Quit()

    @staticmethod
    def _excel_to_pdf(source_path: Path, output_path: Path, client_module: Any) -> None:
        excel = client_module.DispatchEx("Excel.Application")
        excel.Visible = False
        workbook = None
        try:
            workbook = excel.Workbooks.Open(str(source_path))
            workbook.ExportAsFixedFormat(0, str(output_path))
        finally:
            if workbook is not None:
                workbook.Close(False)
            excel.Quit()

    @staticmethod
    def _powerpoint_to_pdf(source_path: Path, output_path: Path, client_module: Any) -> None:
        powerpoint = client_module.DispatchEx("PowerPoint.Application")
        presentation = None
        try:
            presentation = powerpoint.Presentations.Open(str(source_path), WithWindow=False)
            presentation.SaveAs(str(output_path), 32)
        finally:
            if presentation is not None:
                presentation.Close()
            powerpoint.Quit()
